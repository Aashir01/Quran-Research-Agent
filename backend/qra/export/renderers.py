"""Format renderers: Markdown, HTML, Word, slides, PDF.

RTL is the interesting part. Each format handles it differently and none of
them handle it by accident:

* **HTML** — ``dir="rtl"`` plus a Nastaliq font stack; the browser shapes and
  reorders correctly. This is also the PDF source, because a browser is the
  most reliable Arabic typesetter available without bundling fonts and a
  shaping engine.
* **Word / PowerPoint** — the ``w:bidi`` / ``complex script`` properties must be
  set on the paragraph *and* the run, and the font named as a complex-script
  font, or Word lays Arabic out left-to-right in the wrong glyph forms.
* **PDF** — produced by printing the HTML through headless Chromium. Correct
  shaping needs Arabic fonts installed where the rendering happens; when they
  are missing we say so rather than emitting a page of empty boxes.
"""

from __future__ import annotations

import html
import shutil
import subprocess
import tempfile
from pathlib import Path

from qra.export.document import PROVENANCE_LABEL, Document

# Colours match the three provenance states in the web UI, so a printed page
# and a screen say the same thing.
PROVENANCE_COLOUR = {
    "retrieved": "#1f7a4d",
    "system_suggested": "#a4801a",
    "own_note": "#3c5fa8",
}

ARABIC_STACK = "'Amiri Quran','Scheherazade New','Traditional Arabic',serif"
URDU_STACK = "'Noto Nastaliq Urdu','Jameel Noori Nastaleeq','Nafees Nastaleeq',serif"


class ExportUnavailable(RuntimeError):
    """A format needs something this machine does not have. Says what."""


def _label(document: Document, provenance: str) -> str:
    return PROVENANCE_LABEL.get(document.language, PROVENANCE_LABEL["en"]).get(provenance, provenance)


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------


def to_markdown(document: Document) -> str:
    out = [f"# {document.title}"]
    if document.subtitle:
        out.append(f"*{document.subtitle}*")
    out.append("")
    for block in document.blocks:
        if block.kind == "heading":
            out.append(f"{'#' * min(block.level, 6)} {block.text}")
        elif block.kind == "ayah":
            out.append(f"> {block.text}")
            out.append(f"> — **{block.ref}**")
        elif block.kind == "quote":
            out.append(f"> {block.text}")
        elif block.kind == "list":
            out.extend(f"- {item}" for item in block.items)
        elif block.kind == "stat":
            out.append("| | |")
            out.append("|---|---|")
            out.extend(f"| {row[0]} | {row[1]} |" for row in block.rows)
        elif block.kind == "note":
            out.append(f"> [{_label(document, block.provenance)}] {block.text}")
        else:
            out.append(block.text)
        out.append("")
    if document.licence:
        out += ["---", "```", document.licence, "```"]
    out.append(f"\n*{document.generated_at}*")
    return "\n".join(out)


# ---------------------------------------------------------------------------
# HTML (also the PDF source)
# ---------------------------------------------------------------------------


def to_html(document: Document, *, print_ready: bool = True) -> str:
    direction = "rtl" if document.rtl else "ltr"
    parts = [
        "<!doctype html>",
        f'<html lang="{document.language}" dir="{direction}"><head><meta charset="utf-8">',
        f"<title>{html.escape(document.title)}</title><style>",
        f"""
        @page {{ size: A4; margin: 22mm 18mm; }}
        body {{ font-family: system-ui,-apple-system,'Segoe UI',sans-serif; line-height: 1.6;
                color: #14181a; max-width: 800px; margin: 0 auto; padding: 24px; }}
        h1 {{ font-size: 1.7rem; margin-bottom: 2px; }}
        h2 {{ font-size: 1.2rem; margin-top: 26px; border-bottom: 1px solid #dde2df; padding-bottom: 4px; }}
        h3 {{ font-size: 1.02rem; }}
        .subtitle {{ color: #5d6b63; margin-top: 0; }}
        .ayah {{ font-family: {ARABIC_STACK}; font-size: 1.5rem; line-height: 2.2;
                 direction: rtl; text-align: right; margin: 10px 0 2px; }}
        .urdu {{ font-family: {URDU_STACK}; direction: rtl; text-align: right; line-height: 2.4; }}
        .ref {{ color: #5d6b63; font-size: 0.85rem; direction: ltr; text-align: {'left' if document.rtl else 'right'}; }}
        blockquote {{ border-inline-start: 3px solid #c9ced1; margin-inline-start: 0;
                      padding-inline-start: 14px; color: #2c3438; }}
        .prov {{ border-inline-start: 3px solid #c9ced1; padding-inline-start: 12px; margin: 10px 0; }}
        .tag {{ font-size: 0.68rem; text-transform: uppercase; letter-spacing: 0.04em;
                border: 1px solid currentColor; border-radius: 999px; padding: 1px 7px; }}
        table {{ border-collapse: collapse; margin: 10px 0; }}
        td {{ border-bottom: 1px solid #e6eae7; padding: 5px 14px 5px 0; }}
        .licence {{ font-size: 0.78rem; color: #5d6b63; white-space: pre-wrap;
                    border-top: 1px solid #dde2df; margin-top: 28px; padding-top: 10px; }}
        /* Keep an ayah and its reference on the same printed page. */
        .keep {{ break-inside: avoid; }}
        """,
        "</style></head><body>",
        f"<h1>{html.escape(document.title)}</h1>",
    ]
    if document.subtitle:
        parts.append(f'<p class="subtitle">{html.escape(document.subtitle)}</p>')

    for block in document.blocks:
        colour = PROVENANCE_COLOUR.get(block.provenance, "#5d6b63")
        css = "urdu" if block.language == "ur" else ("ayah" if block.kind == "ayah" else "")
        if block.kind == "heading":
            parts.append(f"<h{min(block.level, 4)}>{html.escape(block.text)}</h{min(block.level, 4)}>")
        elif block.kind == "ayah":
            parts.append(
                f'<div class="keep"><p class="ayah">{html.escape(block.text)}</p>'
                f'<p class="ref">{html.escape(block.ref or "")}</p></div>'
            )
        elif block.kind == "quote":
            parts.append(f'<blockquote class="{css}">{html.escape(block.text)}</blockquote>')
        elif block.kind == "list":
            items = "".join(f"<li>{html.escape(i)}</li>" for i in block.items)
            parts.append(f"<ul>{items}</ul>")
        elif block.kind == "stat":
            rows = "".join(
                f"<tr><td>{html.escape(r[0])}</td><td><strong>{html.escape(r[1])}</strong></td></tr>"
                for r in block.rows
            )
            parts.append(f"<table>{rows}</table>")
        elif block.kind == "note":
            parts.append(
                f'<div class="prov" style="border-color:{colour}">'
                f'<span class="tag" style="color:{colour}">{_label(document, block.provenance)}</span>'
                f'<p class="{css}">{html.escape(block.text)}</p></div>'
            )
        else:
            parts.append(f'<p class="{css}">{html.escape(block.text)}</p>')

    if document.licence:
        parts.append(f'<div class="licence">{html.escape(document.licence)}</div>')
    parts.append(f'<p class="ref">{document.generated_at}</p>')
    parts.append("</body></html>")
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------


def _set_rtl(paragraph, font_name: str) -> None:
    """Word needs bidi on the paragraph and complex-script settings on each run.

    Without ``w:rtl`` on the run and ``w:cs`` on the font, Word renders Arabic
    with Latin layout rules: correct characters, wrong order and wrong joining.
    """
    from docx.oxml.ns import qn
    from docx.shared import Pt

    p_pr = paragraph._p.get_or_add_pPr()
    bidi = p_pr.makeelement(qn("w:bidi"), {})
    p_pr.append(bidi)
    for run in paragraph.runs:
        r_pr = run._r.get_or_add_rPr()
        r_pr.append(r_pr.makeelement(qn("w:rtl"), {}))
        r_pr.append(r_pr.makeelement(qn("w:cs"), {}))
        run.font.name = font_name
        r_fonts = r_pr.find(qn("w:rFonts"))
        if r_fonts is None:
            r_fonts = r_pr.makeelement(qn("w:rFonts"), {})
            r_pr.append(r_fonts)
        r_fonts.set(qn("w:cs"), font_name)
        run.font.size = run.font.size or Pt(16)


def to_docx(document: Document, path: Path) -> Path:
    from docx import Document as Docx
    from docx.shared import Pt, RGBColor

    docx = Docx()
    heading = docx.add_heading(document.title, level=1)
    if document.rtl:
        _set_rtl(heading, "Noto Nastaliq Urdu")
    if document.subtitle:
        docx.add_paragraph(document.subtitle)

    for block in document.blocks:
        if block.kind == "heading":
            docx.add_heading(block.text, level=min(block.level, 4))
        elif block.kind == "ayah":
            paragraph = docx.add_paragraph()
            run = paragraph.add_run(block.text)
            run.font.size = Pt(18)
            _set_rtl(paragraph, "Amiri Quran")
            reference = docx.add_paragraph()
            ref_run = reference.add_run(block.ref or "")
            ref_run.font.size = Pt(9)
            ref_run.font.color.rgb = RGBColor(0x5D, 0x6B, 0x63)
        elif block.kind == "list":
            for item in block.items:
                docx.add_paragraph(item, style="List Bullet")
        elif block.kind == "stat":
            table = docx.add_table(rows=0, cols=2)
            table.style = "Light Grid Accent 1"
            for row in block.rows:
                cells = table.add_row().cells
                cells[0].text, cells[1].text = row[0], row[1]
        elif block.kind in ("quote", "note"):
            paragraph = docx.add_paragraph()
            tag = paragraph.add_run(f"[{_label(document, block.provenance)}] ")
            tag.bold = True
            colour = PROVENANCE_COLOUR.get(block.provenance, "#5d6b63").lstrip("#")
            tag.font.color.rgb = RGBColor(*(int(colour[i : i + 2], 16) for i in (0, 2, 4)))
            paragraph.add_run(block.text)
            if block.rtl:
                _set_rtl(paragraph, "Noto Nastaliq Urdu" if block.language == "ur" else "Amiri Quran")
        else:
            paragraph = docx.add_paragraph(block.text)
            if block.rtl:
                _set_rtl(paragraph, "Noto Nastaliq Urdu" if block.language == "ur" else "Amiri Quran")

    if document.licence:
        docx.add_page_break()
        docx.add_paragraph(document.licence).runs[0].font.size = Pt(8)
    docx.add_paragraph(document.generated_at).runs[0].font.size = Pt(8)

    docx.save(path)
    return path


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def to_pptx(document: Document, path: Path, *, max_lines_per_slide: int = 6) -> Path:
    """One slide per heading, with its following blocks as bullets.

    Ayat get a slide to themselves at a readable size — a verse squeezed into a
    bullet list is unreadable in Nastaliq or Uthmani at projector distance.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    blank, title_layout = presentation.slide_layouts[6], presentation.slide_layouts[0]

    opening = presentation.slides.add_slide(title_layout)
    opening.shapes.title.text = document.title
    if document.subtitle:
        opening.placeholders[1].text = document.subtitle

    current = None
    lines = 0

    def new_slide(heading: str) -> tuple:
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(1.0))
        frame = box.text_frame
        frame.text = heading
        frame.paragraphs[0].runs[0].font.size = Pt(28)
        frame.paragraphs[0].runs[0].font.bold = True
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.2))
        body.text_frame.word_wrap = True
        return slide, body.text_frame

    for block in document.blocks:
        if block.kind == "heading":
            current = new_slide(block.text)[1]
            lines = 0
            continue
        if current is None:
            current = new_slide(document.title)[1]
            lines = 0

        if block.kind == "ayah":
            slide = presentation.slides.add_slide(blank)
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(3.4))
            frame = box.text_frame
            frame.word_wrap = True
            frame.text = block.text
            paragraph = frame.paragraphs[0]
            paragraph.runs[0].font.size = Pt(30)
            paragraph.runs[0].font.name = "Amiri Quran"
            _rtl_pptx(paragraph)
            ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9.0), Inches(0.6))
            ref_box.text_frame.text = block.ref or ""
            ref_box.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            current, lines = None, 0
            continue

        texts = block.items if block.kind == "list" else [block.text]
        for text in texts:
            if lines >= max_lines_per_slide:
                current = new_slide("(continued)")[1]
                lines = 0
            paragraph = current.add_paragraph()
            paragraph.text = f"• {text[:220]}"
            paragraph.runs[0].font.size = Pt(16)
            if block.rtl:
                paragraph.runs[0].font.name = "Noto Nastaliq Urdu"
                _rtl_pptx(paragraph)
            lines += 1

    presentation.save(path)
    return path


def _rtl_pptx(paragraph) -> None:
    """Set right-to-left on a PowerPoint paragraph (python-pptx has no API)."""
    paragraph._pPr.set("rtl", "1") if paragraph._pPr is not None else None
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set("rtl", "1")
    p_pr.set("algn", "r")


# ---------------------------------------------------------------------------
# PDF — printed from the HTML by headless Chromium
# ---------------------------------------------------------------------------


def _chromium() -> str | None:
    for candidate in (
        "/opt/pw-browsers/chromium/chrome-linux/chrome",
        "/opt/pw-browsers/chromium-1194/chrome-linux/chrome",
    ):
        if Path(candidate).exists():
            return candidate
    return shutil.which("chromium") or shutil.which("chromium-browser") or shutil.which("google-chrome")


# Fonts that actually typeset Qur'anic Arabic or Urdu properly. DejaVu and
# Unifont answer `fc-list :lang=ar` because they carry *some* Arabic glyphs,
# but they do not shape or join it acceptably — treating them as sufficient
# would ship a PDF that is technically not tofu and still unusable.
_QUALITY_ARABIC_FONTS = (
    "amiri", "scheherazade", "noto naskh", "noto nastaliq", "traditional arabic",
    "kfgqpc", "lateef", "jameel", "nafees", "me_quran",
)


def arabic_font_status() -> str:
    """``quality`` | ``fallback`` | ``none``."""
    if not shutil.which("fc-list"):
        return "none"
    out = subprocess.run(["fc-list", ":lang=ar"], capture_output=True, text=True, check=False)
    listing = out.stdout.lower()
    if not listing.strip():
        return "none"
    return "quality" if any(f in listing for f in _QUALITY_ARABIC_FONTS) else "fallback"


def to_pdf(document: Document, path: Path) -> Path:  # noqa: C901
    """Print the HTML export through headless Chromium.

    A browser is used rather than a PDF library because it already does Arabic
    shaping and bidi correctly. It still needs fonts: if none covering Arabic
    are installed, we refuse rather than produce a document full of empty boxes,
    and point at the HTML export, which renders correctly in the researcher's
    own browser.
    """
    binary = _chromium()
    if binary is None:
        raise ExportUnavailable(
            "PDF export needs headless Chromium on this machine. The HTML export is "
            "print-ready — open it in a browser and print to PDF."
        )
    status = arabic_font_status()
    if status == "none":
        raise ExportUnavailable(
            "No Arabic-covering fonts are installed here, so a PDF would render scripture as "
            "empty boxes. Install e.g. fonts-noto-naskh-arabic and fonts-noto-nastaliq-urdu, "
            "or use the HTML export and print from a machine that has them."
        )

    with tempfile.TemporaryDirectory() as tmp:
        source = Path(tmp) / "document.html"
        source.write_text(to_html(document), encoding="utf-8")
        subprocess.run(
            [
                binary,
                "--headless",
                "--disable-gpu",
                "--no-sandbox",
                "--no-pdf-header-footer",
                f"--print-to-pdf={path}",
                source.as_uri(),
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
    if status == "fallback":
        # Produced, but the reader should know the typography is not a proper
        # Qur'anic face — silently shipping DejaVu-shaped scripture would be
        # the wrong kind of quiet.
        path.with_suffix(".fonts.txt").write_text(
            "This PDF was rendered without a Qur'anic Arabic typeface installed. The text is "
            "correct but the shaping is a generic fallback. Install fonts-noto-naskh-arabic and "
            "fonts-noto-nastaliq-urdu, or print the HTML export from a machine that has them.\n",
            encoding="utf-8",
        )
    return path


RENDERERS = {
    "md": to_markdown,
    "html": to_html,
    "docx": to_docx,
    "pptx": to_pptx,
    "pdf": to_pdf,
}
